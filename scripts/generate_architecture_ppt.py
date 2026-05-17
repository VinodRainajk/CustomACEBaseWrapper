"""
Generates docs/CustomACEBaseWrapper-Architecture-Overview.pptx
Run: python scripts/generate_architecture_ppt.py
Requires: pip install python-pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "CustomACEBaseWrapper-Architecture-Overview.pptx"


def add_title_slide(prs, title: str, subtitle: str):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    ph = slide.placeholders[1]
    ph.text = subtitle
    return slide


def add_bullet_slide(prs, title: str, bullets: list[str]):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
    return slide


def add_two_column_bullets(prs, title: str, left_title: str, left: list[str], right_title: str, right: list[str]):
    """Title + two text boxes using blank layout."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Title
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tf = tbox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    col_w = Inches(4.4)
    y = Inches(1.35)
    h = Inches(5.2)

    for x, sub, items in (
        (Inches(0.5), left_title, left),
        (Inches(5.1), right_title, right),
    ):
        box = slide.shapes.add_textbox(x, y, col_w, h)
        btf = box.text_frame
        btf.word_wrap = True
        p0 = btf.paragraphs[0]
        p0.text = sub
        p0.font.size = Pt(20)
        p0.font.bold = True
        p0.font.color.rgb = RGBColor(0x2E, 0x5E, 0xAA)
        for line in items:
            p = btf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(15)
            p.space_before = Pt(4)
    return slide


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Custom ACE Base Wrapper",
        "Architecture overview\nJava 17 · Maven library · Cucumber BDD · ace-base integration\nMay 2026",
    )

    add_bullet_slide(
        prs,
        "What is this project?",
        [
            "Maven artifact: com.qa.framework — custom-ace-base-wrapper (JAR)",
            "Reusable QA library: Cucumber step definitions, hooks, runners, config helpers",
            "Pairs with consumer projects (features, YAML profiles, payloads on their classpath)",
            "Integrates ace-base for TestNG-driven UI/API runs (glue + TestNGRunner lifecycle)",
        ],
    )

    add_bullet_slide(
        prs,
        "What the library contains vs. the test project",
        [
            "In the JAR: step definitions (db, api, ui, payload), hooks, DB/API utilities, runners",
            "In the test project: .feature files, config/{profile}/*.yaml, payloads, SQL assets",
            "Install locally: mvn clean install — then depend on artifact 1.0.0 from test POMs",
        ],
    )

    add_bullet_slide(
        prs,
        "High-level architecture (layers)",
        [
            "Consumer: Gherkin features + YAML (master + feature override) + JSON/YAML payloads",
            "Runners: JUnit Platform (DB) or TestNG + ace-base (UI/API)",
            "Glue: com.qa.framework.stepdefinitions.{db,api,ui} + com.qa.framework.payload",
            "Config: UnifiedConfigLoader merges config/{profile}/master-config.yaml + {feature}-config.yaml",
            "Runtime: DatabaseManager / JDBC; RestAssured for API; payload loaders",
            "Foundation: BaseWrapper, ConfigurableWrapper, WrapperException; shared utils",
        ],
    )

    add_two_column_bullets(
        prs,
        "Execution paths",
        "Database (JUnit)",
        [
            "InterceptorRunner",
            "  → DBTestRunner",
            "Tag filter: @DB",
            "Glue: stepdefinitions.db, payload",
            "DatabaseHooks + DatabaseConfigLoader",
            "Closes connections in @After",
        ],
        "UI / API (TestNG)",
        [
            "UIAPITestNGRunner",
            "extends ace-base TestNGRunner",
            "Glue: com.acebase.glue +",
            "  stepdefinitions.ui, api, payload",
            "APIHooks loads APIConfig per feature",
            "ace-base CucumberPlugin + reports",
        ],
    )

    add_bullet_slide(
        prs,
        "Configuration model",
        [
            "Active profile: -Dprofile=local|dev|qa|staging|prod (default: local)",
            "Unified YAML under classpath: config/{profile}/",
            "Merge order: master-config.yaml, then optional {feature}-config.yaml (later wins)",
            "Database: top-level db block; optional sections.{scenarioName}.db for overrides",
            "API: APIConfigLoader / APIConfig aligned with merged YAML (see API_ARCHITECTURE.md)",
        ],
    )

    add_bullet_slide(
        prs,
        "Cucumber glue (main packages)",
        [
            "stepdefinitions.db — connection, CRUD, procedures, validation, file/SQL steps; DatabaseHooks",
            "stepdefinitions.api — request, status, body, configuration, errors; APIHooks + APIStepContext",
            "stepdefinitions.ui — UI scenarios (extends consumer/driver setup via ace-base where used)",
            "payload — Feature payload / wait-related steps; works with classpath payload resources",
            "util.EnvTagFilter — builds tag expressions: type tag AND env (@all, @dev, @prod, @nonProd)",
        ],
    )

    add_bullet_slide(
        prs,
        "Key runtime components",
        [
            "DB: DatabaseConnectionFactory, DatabaseConnection, DatabaseManager, PendingStatementExecutor",
            "API: RestAssured; APIStepContext holds base URL, config, last response",
            "Utils: DynamicValueUtils, PollingUtils, CsvMismatchReportWriter, FileSearchUtils, RegexUtils, StringUtils",
        ],
    )

    add_bullet_slide(
        prs,
        "Dependencies (pom.xml snapshot)",
        [
            "Cucumber 7.x, JUnit Platform suite, cucumber-testng, TestNG 7.8",
            "ace-base 1.0.0 — TestNG runner, shared glue",
            "JDBC: MySQL, PostgreSQL, SQL Server, Oracle drivers",
            "RestAssured 5.4, json-path, SnakeYAML 2.2, Commons IO",
        ],
    )

    add_bullet_slide(
        prs,
        "Documentation in repo",
        [
            "README.md — library usage and install",
            "ARCHITECTURE.md — DB-centric flow (some class names evolved; cross-check code)",
            "API_ARCHITECTURE.md — API steps and config/payload conventions",
            "CONFIGURATION.md, DATABASE_CONFIGURATION.md, FEATURE_PAYLOAD.md, QUICK_START.md",
        ],
    )

    add_title_slide(
        prs,
        "Thank you",
        "Regenerate this deck: python scripts/generate_architecture_ppt.py\nOutput: docs/CustomACEBaseWrapper-Architecture-Overview.pptx",
    )

    prs.save(OUT)
    print(f"Wrote: {OUT}")


if __name__ == "__main__":
    main()
